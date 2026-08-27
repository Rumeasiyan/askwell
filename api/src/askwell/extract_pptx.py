"""PowerPoint extraction. `M1-EXTRACT-ING-027`.

**`python-pptx`, MIT-licensed.** One slide, one anchor — a deck has a real
ordinal a PDF or a Word file does not, so `page_number` here needs no
"approximate" hedge the way `extract_docx` does.

**Speaker notes are included and labelled**, per the ticket's own edge case,
rather than dropped or left indistinguishable from the slide's own text — a
retrieval hit on a note a presenter wrote for themselves should read as one.
"""

import asyncio
from typing import TYPE_CHECKING, cast

from pptx import Presentation as open_presentation
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.table import Table
from pptx.text.text import TextFrame

from askwell.extract_common import Anchor, write_anchors
from askwell.logging import get_logger

if TYPE_CHECKING:
    from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker

    from askwell.ingest import Report, Work

log = get_logger(__name__)

ANCHOR_KIND = "slide"


def _shape_text(shape: BaseShape) -> str:
    if shape.has_table:
        table = cast(Table, shape.table)  # type: ignore[attr-defined]
        lines = ["[TABLE]"]
        for row in table.rows:
            lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        lines.append("[/TABLE]")
        return "\n".join(lines)
    if shape.has_text_frame:
        frame = cast(TextFrame, shape.text_frame)  # type: ignore[attr-defined]
        lines = []
        for paragraph in frame.paragraphs:
            body = paragraph.text.strip()
            if body:
                lines.append(f"- {body}" if paragraph.level else body)
        return "\n".join(lines)
    return ""


def _slide_text(slide: Slide) -> str:
    parts = [piece for shape in slide.shapes if (piece := _shape_text(shape))]

    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        notes = notes_frame.text.strip() if notes_frame is not None else ""
        if notes:
            parts.append(f"[Speaker notes]\n{notes}")

    return "\n\n".join(parts)


def _load(path: str) -> Presentation:
    return open_presentation(path)


async def run(work: "Work", report: "Report", factory: "async_sessionmaker[AsyncSession]") -> None:
    presentation = await asyncio.to_thread(_load, work.path)
    slides = list(presentation.slides)

    anchors: list[Anchor] = []
    for index, slide in enumerate(slides, start=1):
        body = _slide_text(slide).strip()
        anchors.append(
            Anchor(
                page_number=index,
                label=f"Slide {index}",
                text=body or None,
                has_text=bool(body),
            )
        )
        await report(index, len(slides))

    await write_anchors(factory, work, anchors, ANCHOR_KIND)

    log.info(
        "extract_pptx_completed",
        document_id=str(work.document_id),
        filename=work.filename,
        slides=len(slides),
    )
