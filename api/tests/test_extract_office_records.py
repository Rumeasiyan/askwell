"""Word, PowerPoint, Excel, text, Markdown and HTML extraction, end to end,
against a real Postgres. `M1-EXTRACT-ING-027`.

Mirrors `test_ingest_records.py`'s pattern for `extract_pdf`: files are put
through the real `add()` path so `askwell.filetypes.detect` decides the mime
type from actual bytes, then `ingest.process` runs the real `extract` stage
dispatcher — nothing here is a stand-in extractor. Six documents, one per
new format, each asserted for the thing that format's edge case in the ticket
actually asks for: headings and tables for Word, labelled speaker notes for
PowerPoint, a sheet-and-row anchor for Excel, front matter excluded from
Markdown, navigation chrome discarded from HTML — plus the validation rule
that a document with nothing extractable in it fails with a reason rather
than reaching `ready` empty.
"""

import uuid
from collections.abc import AsyncIterator, Iterator
from pathlib import Path

import docx
import openpyxl
import pytest
import pytest_asyncio
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

from askwell import ingest
from askwell.config import Settings

from .test_ingest_records import TABLES, nominate, recorded

pytestmark = pytest.mark.requires_db


# Duplicated from `test_ingest_records.py` rather than imported: a fixture
# reused across test modules by import is flagged by ruff (F811) the moment a
# test function's own parameter shadows the imported name, which every test
# below does. `TABLES`, `nominate` and `recorded` are plain functions/data and
# have no such conflict, so those are shared for real.


@pytest.fixture
def async_url(database_url: str) -> str:
    return database_url.replace("postgresql://", "postgresql+psycopg://", 1)


@pytest_asyncio.fixture
async def factory(async_url: str) -> AsyncIterator[async_sessionmaker[AsyncSession]]:
    engine = create_async_engine(async_url)
    sessions = async_sessionmaker(engine, expire_on_commit=False)
    async with sessions() as opened:
        await opened.execute(text(f"TRUNCATE {TABLES} CASCADE"))
        await opened.commit()
    yield sessions
    async with sessions() as opened:
        await opened.execute(text(f"TRUNCATE {TABLES} CASCADE"))
        await opened.commit()
    await engine.dispose()


@pytest_asyncio.fixture
async def session(
    factory: async_sessionmaker[AsyncSession],
) -> AsyncIterator[AsyncSession]:
    async with factory() as opened:
        yield opened
        await opened.rollback()


@pytest.fixture
def unreachable_queue(settings: Settings, monkeypatch: pytest.MonkeyPatch) -> Iterator[Settings]:
    sent: list[uuid.UUID] = []

    async def fake_dispatch(
        _settings: Settings,
        document_ids: list[uuid.UUID],
        **_kwargs: object,
    ) -> int:
        sent.extend(document_ids)
        return len(document_ids)

    monkeypatch.setattr(ingest, "dispatch", fake_dispatch)
    yield settings


def _write_docx(path: Path) -> None:
    document = docx.Document()
    document.add_heading("Renewal Terms", level=1)
    document.add_paragraph("Either party may terminate on ninety days written notice.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Item"
    table.rows[0].cells[1].text = "Price"
    table.rows[1].cells[0].text = "Widget"
    table.rows[1].cells[1].text = "9.99"
    document.save(str(path))


def _write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Quarterly results are ahead of forecast."
    slide.notes_slide.notes_text_frame.text = "Mention the pricing change."
    presentation.save(str(path))


def _write_xlsx(path: Path) -> None:
    workbook = openpyxl.Workbook()
    budget = workbook.active
    budget.title = "Budget"
    budget.append(["Item", "Cost"])
    budget.append(["Widgets", 42])
    other = workbook.create_sheet("Notes")
    other.append(["Remember to renew the lease."])
    workbook.save(str(path))


def _write_empty_xlsx(path: Path) -> None:
    workbook = openpyxl.Workbook()
    workbook.save(str(path))


async def _process(
    factory: async_sessionmaker[AsyncSession],
    settings: Settings,
    session: AsyncSession,
    tmp_path: Path,
    filename: str,
) -> tuple[str, list[tuple], object]:
    await nominate(session, str(tmp_path))
    documents = await recorded(session, tmp_path, filename)
    document_id = documents[0]
    outcome = await ingest.process(factory, settings, document_id)
    pages = (
        await session.execute(
            text(
                "SELECT page_number, text, has_text, anchor_label FROM document_pages "
                "WHERE document_id = :id ORDER BY page_number"
            ),
            {"id": document_id},
        )
    ).all()
    return outcome, [tuple(row) for row in pages], document_id


async def test_word_headings_and_tables_survive_as_structural_markers(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    _write_docx(tmp_path / "contract.docx")
    outcome, pages, document_id = await _process(
        factory, unreachable_queue, session, tmp_path, "contract.docx"
    )

    assert outcome == "parked"
    assert len(pages) == 1
    body = pages[0][1]
    assert "# Renewal Terms" in body
    assert "Either party may terminate" in body
    assert "[TABLE]" in body and "Widget | 9.99" in body

    anchor_kind = (
        await session.execute(
            text("SELECT anchor_kind FROM documents WHERE id = :id"),
            {"id": document_id},
        )
    ).scalar_one()
    assert anchor_kind == "page"


async def test_powerpoint_speaker_notes_are_included_and_labelled(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    _write_pptx(tmp_path / "deck.pptx")
    outcome, pages, _ = await _process(factory, unreachable_queue, session, tmp_path, "deck.pptx")

    assert outcome == "parked"
    assert len(pages) == 1
    assert pages[0][3] == "Slide 1"
    assert "Quarterly results are ahead of forecast." in pages[0][1]
    assert "[Speaker notes]" in pages[0][1]
    assert "Mention the pricing change." in pages[0][1]


async def test_excel_is_read_as_one_anchor_per_row_across_sheets(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    _write_xlsx(tmp_path / "figures.xlsx")
    outcome, pages, _ = await _process(
        factory, unreachable_queue, session, tmp_path, "figures.xlsx"
    )

    assert outcome == "parked"
    labels = [row[3] for row in pages]
    assert "Budget, row 1" in labels
    assert "Budget, row 2" in labels
    assert "Notes, row 1" in labels
    widget_row = next(row for row in pages if row[3] == "Budget, row 2")
    assert widget_row[1] == "Widgets | 42"


async def test_markdown_front_matter_is_excluded_from_the_first_section(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    (tmp_path / "note.md").write_text(
        "---\ntitle: Renewal notice\n---\n"
        "# Renewal\n\nEither party may terminate on ninety days written notice.\n"
    )
    outcome, pages, _ = await _process(factory, unreachable_queue, session, tmp_path, "note.md")

    assert outcome == "parked"
    all_text = "\n".join(row[1] or "" for row in pages)
    assert "title: Renewal notice" not in all_text
    assert "Either party may terminate" in all_text
    assert any(row[3] == "Renewal" for row in pages)


async def test_html_navigation_chrome_does_not_reach_extracted_text(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    (tmp_path / "page.html").write_text(
        "<html><head><title>Contract</title></head><body>"
        "<nav><a href='/'>Home</a><a href='/about'>About</a></nav>"
        "<h1>Terms</h1><p>Either party may terminate on ninety days notice.</p>"
        "</body></html>"
    )
    outcome, pages, _ = await _process(factory, unreachable_queue, session, tmp_path, "page.html")

    assert outcome == "parked"
    all_text = "\n".join(row[1] or "" for row in pages)
    assert "Home" not in all_text
    assert "About" not in all_text
    assert "Contract" not in all_text  # the <title>, not the page's content
    assert "Either party may terminate" in all_text


async def test_plain_text_with_no_headings_is_one_section(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    (tmp_path / "note.txt").write_text("Just a plain note, no structure at all.")
    outcome, pages, _ = await _process(factory, unreachable_queue, session, tmp_path, "note.txt")

    assert outcome == "parked"
    assert len(pages) == 1
    assert pages[0][3] is None
    assert pages[0][1] == "Just a plain note, no structure at all."


async def test_a_document_with_no_extractable_text_fails_with_a_reason(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    """The ticket's own validation rule: never an empty indexed document."""
    _write_empty_xlsx(tmp_path / "blank.xlsx")
    outcome, _, document_id = await _process(
        factory, unreachable_queue, session, tmp_path, "blank.xlsx"
    )

    assert outcome == "failed"

    row = (
        await session.execute(
            text("SELECT status FROM documents WHERE id = :id"), {"id": document_id}
        )
    ).scalar_one()
    assert row != "ready"

    error = (
        await session.execute(
            text("SELECT error FROM ingest_jobs WHERE document_id = :id"), {"id": document_id}
        )
    ).scalar_one()
    assert "could not find any text" in error


async def test_a_legacy_binary_word_file_fails_by_name_rather_than_crashing(
    factory: async_sessionmaker[AsyncSession],
    session: AsyncSession,
    tmp_path: Path,
    unreachable_queue: Settings,
) -> None:
    (tmp_path / "old.doc").write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 32)
    await nominate(session, str(tmp_path))
    documents = await recorded(session, tmp_path, "old.doc")

    outcome = await ingest.process(factory, unreachable_queue, documents[0])

    assert outcome == "failed"
    error = (
        await session.execute(
            text("SELECT error FROM ingest_jobs WHERE document_id = :id"), {"id": documents[0]}
        )
    ).scalar_one()
    assert "does not read the older binary Office formats" in error
