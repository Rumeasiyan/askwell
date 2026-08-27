"""PowerPoint extraction's own rules, without a database. `M1-EXTRACT-ING-027`."""

from pptx import Presentation
from pptx.util import Inches

from askwell.extract_pptx import _shape_text, _slide_text


def _blank_slide(presentation: Presentation) -> object:
    layout = presentation.slide_layouts[6]  # the blank layout
    return presentation.slides.add_slide(layout)


def test_a_textbox_shape_becomes_its_paragraph_text() -> None:
    presentation = Presentation()
    slide = _blank_slide(presentation)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Quarterly results"
    assert _shape_text(box) == "Quarterly results"


def test_a_table_shape_becomes_bracketed_rows() -> None:
    presentation = Presentation()
    slide = _blank_slide(presentation)
    graphic_frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    table = graphic_frame.table
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Price"
    table.cell(1, 0).text = "Widget"
    table.cell(1, 1).text = "9.99"
    assert _shape_text(graphic_frame).splitlines() == [
        "[TABLE]",
        "Item | Price",
        "Widget | 9.99",
        "[/TABLE]",
    ]


def test_speaker_notes_are_included_and_labelled() -> None:
    presentation = Presentation()
    slide = _blank_slide(presentation)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "On screen text."
    slide.notes_slide.notes_text_frame.text = "Remember to mention the pricing change."

    body = _slide_text(slide)
    assert "On screen text." in body
    assert "[Speaker notes]" in body
    assert "Remember to mention the pricing change." in body


def test_a_slide_with_no_notes_has_no_notes_label() -> None:
    presentation = Presentation()
    slide = _blank_slide(presentation)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Just the slide."
    assert "[Speaker notes]" not in _slide_text(slide)
